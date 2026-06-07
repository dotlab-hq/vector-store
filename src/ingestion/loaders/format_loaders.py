"""
Format-specific document loaders built on top of LangChain document loaders.

Each loader handles a single format family (PDF, DOCX, PPTX, HTML, CSV, JSON, Excel)
and converts the output into our RawDocument format for consistent pipeline handling.
"""

from __future__ import annotations

import csv
import io
import json
from pathlib import Path

import structlog

from src.ingestion.loaders.base import DocumentLoader, RawDocument

logger = structlog.get_logger()


# ---------------------------------------------------------------------------
# PDF Loader (pypdf — text PDFs)
# ---------------------------------------------------------------------------
class PDFLoader(DocumentLoader):
    """PDF loader via pypdf for text-based PDFs."""

    def can_load(self, file_path: Path) -> bool:
        return file_path.suffix.lower() == ".pdf"

    def load(self, file_path: Path) -> RawDocument:
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(file_path))
            pages: list[str] = []
            for page in reader.pages:
                text = page.extract_text()
                if text and text.strip():
                    pages.append(text.strip())

            combined = "\n\n".join(pages)

            metadata: dict = {
                "title": file_path.stem,
                "extension": ".pdf",
                "converter": "pypdf",
                "page_count": len(reader.pages),
            }

            return RawDocument(
                content=combined,
                metadata=metadata,
                source_path=str(file_path),
            )
        except ImportError:
            logger.warning("pypdf_not_installed", file=str(file_path))
            return RawDocument(content="", metadata={}, source_path=str(file_path))
        except Exception as e:
            logger.warning("pdf_load_failed", file=str(file_path), error=str(e))
            return RawDocument(content="", metadata={}, source_path=str(file_path))


# ---------------------------------------------------------------------------
# DOCX Loader
# ---------------------------------------------------------------------------
class DOCXLoader(DocumentLoader):
    """Word document loader via python-docx."""

    def can_load(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in {".docx", ".doc"}

    def load(self, file_path: Path) -> RawDocument:
        try:
            from docx import Document

            doc = Document(str(file_path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            content = "\n\n".join(paragraphs)

            metadata: dict = {
                "title": file_path.stem,
                "extension": file_path.suffix.lower(),
                "converter": "python-docx",
                "paragraph_count": len(paragraphs),
            }
            if doc.core_properties.title:
                metadata["doc_title"] = doc.core_properties.title
            if doc.core_properties.author:
                metadata["author"] = doc.core_properties.author

            return RawDocument(
                content=content,
                metadata=metadata,
                source_path=str(file_path),
            )
        except ImportError:
            logger.warning("python_docx_not_installed", file=str(file_path))
            return RawDocument(content="", metadata={}, source_path=str(file_path))
        except Exception as e:
            logger.warning("docx_load_failed", file=str(file_path), error=str(e))
            return RawDocument(content="", metadata={}, source_path=str(file_path))


# ---------------------------------------------------------------------------
# PPTX Loader
# ---------------------------------------------------------------------------
class PPTXLoader(DocumentLoader):
    """PowerPoint presentation loader."""

    def can_load(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in {".pptx", ".ppt"}

    def load(self, file_path: Path) -> RawDocument:
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))
            slide_texts: list[str] = []
            for i, slide in enumerate(prs.slides):
                parts: list[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
                if parts:
                    slide_texts.append(f"# Slide {i + 1}\n" + "\n".join(parts))

            content = "\n\n".join(slide_texts)

            metadata: dict = {
                "title": file_path.stem,
                "extension": file_path.suffix.lower(),
                "converter": "python-pptx",
                "slide_count": len(prs.slides),
            }

            return RawDocument(
                content=content,
                metadata=metadata,
                source_path=str(file_path),
            )
        except ImportError:
            logger.warning("python_pptx_not_installed", file=str(file_path))
            return RawDocument(content="", metadata={}, source_path=str(file_path))
        except Exception as e:
            logger.warning("pptx_load_failed", file=str(file_path), error=str(e))
            return RawDocument(content="", metadata={}, source_path=str(file_path))


# ---------------------------------------------------------------------------
# HTML Loader
# ---------------------------------------------------------------------------
class HTMLLoader(DocumentLoader):
    """HTML page loader via BeautifulSoup."""

    def can_load(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in {".html", ".htm"}

    def load(self, file_path: Path) -> RawDocument:
        try:
            from bs4 import BeautifulSoup

            raw_html = file_path.read_text(encoding="utf-8", errors="replace")
            soup = BeautifulSoup(raw_html, "lxml")

            # Remove script/style tags
            for tag in soup(["script", "style", "noscript"]):
                tag.decompose()

            title = ""
            if soup.title and soup.title.string:
                title = soup.title.string.strip()

            content = soup.get_text(separator="\n", strip=True)

            metadata: dict = {
                "title": title or file_path.stem,
                "extension": file_path.suffix.lower(),
                "converter": "beautifulsoup",
            }

            return RawDocument(
                content=content,
                metadata=metadata,
                source_path=str(file_path),
            )
        except ImportError:
            logger.warning("beautifulsoup_not_installed", file=str(file_path))
            return RawDocument(content="", metadata={}, source_path=str(file_path))
        except Exception as e:
            logger.warning("html_load_failed", file=str(file_path), error=str(e))
            return RawDocument(content="", metadata={}, source_path=str(file_path))


# ---------------------------------------------------------------------------
# CSV Loader
# ---------------------------------------------------------------------------
class CSVLoader(DocumentLoader):
    """CSV file loader — preserves all rows as structured text."""

    def can_load(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in {".csv", ".tsv"}

    def load(self, file_path: Path) -> RawDocument:
        delimiter = "\t" if file_path.suffix.lower() == ".tsv" else ","
        try:
            raw_text = file_path.read_text(encoding="utf-8", errors="replace")
            reader = csv.DictReader(io.StringIO(raw_text), delimiter=delimiter)
            rows = list(reader)
            if not rows:
                return RawDocument(
                    content="",
                    metadata={"title": file_path.stem, "extension": file_path.suffix.lower(), "converter": "csv"},
                    source_path=str(file_path),
                )

            headers = list(rows[0].keys())
            lines = [" | ".join(headers)]
            for row in rows:
                lines.append(" | ".join(str(row.get(h, "")) for h in headers))

            content = "\n".join(lines)

            metadata: dict = {
                "title": file_path.stem,
                "extension": file_path.suffix.lower(),
                "converter": "csv",
                "row_count": len(rows),
                "columns": headers,
            }

            return RawDocument(content=content, metadata=metadata, source_path=str(file_path))
        except Exception as e:
            logger.warning("csv_load_failed", file=str(file_path), error=str(e))
            return RawDocument(content="", metadata={}, source_path=str(file_path))


# ---------------------------------------------------------------------------
# JSON Loader
# ---------------------------------------------------------------------------
class JSONLoader(DocumentLoader):
    """JSON / JSON Lines file loader."""

    def can_load(self, file_path: Path) -> bool:
        return file_path.suffix.lower() == ".json"

    def load(self, file_path: Path) -> RawDocument:
        try:
            raw_text = file_path.read_text(encoding="utf-8", errors="replace")
            data = json.loads(raw_text)

            if isinstance(data, list):
                content = json.dumps(data, indent=2, ensure_ascii=False)
                record_count = len(data)
            elif isinstance(data, dict):
                content = json.dumps(data, indent=2, ensure_ascii=False)
                record_count = 1
            else:
                content = str(data)
                record_count = 0

            metadata: dict = {
                "title": file_path.stem,
                "extension": ".json",
                "converter": "json",
                "record_count": record_count,
            }

            return RawDocument(content=content, metadata=metadata, source_path=str(file_path))
        except Exception as e:
            logger.warning("json_load_failed", file=str(file_path), error=str(e))
            return RawDocument(content="", metadata={}, source_path=str(file_path))


# ---------------------------------------------------------------------------
# Excel Loader
# ---------------------------------------------------------------------------
class ExcelLoader(DocumentLoader):
    """Excel spreadsheet loader via openpyxl."""

    def can_load(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in {".xlsx", ".xls", ".ods"}

    def load(self, file_path: Path) -> RawDocument:
        try:
            from openpyxl import load_workbook

            wb = load_workbook(str(file_path), read_only=True, data_only=True)
            sheet_parts: list[str] = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = list(ws.iter_rows(values_only=True))
                if not rows:
                    continue

                headers = [str(cell) if cell is not None else "" for cell in rows[0]]
                lines = [f"# Sheet: {sheet_name}", " | ".join(headers)]
                for row in rows[1:]:
                    lines.append(" | ".join(str(cell) if cell is not None else "" for cell in row))
                sheet_parts.append("\n".join(lines))

            wb.close()
            content = "\n\n".join(sheet_parts)

            metadata: dict = {
                "title": file_path.stem,
                "extension": file_path.suffix.lower(),
                "converter": "openpyxl",
                "sheet_count": len(wb.sheetnames),
                "sheet_names": wb.sheetnames,
            }

            return RawDocument(content=content, metadata=metadata, source_path=str(file_path))
        except ImportError:
            logger.warning("openpyxl_not_installed", file=str(file_path))
            return RawDocument(content="", metadata={}, source_path=str(file_path))
        except Exception as e:
            logger.warning("excel_load_failed", file=str(file_path), error=str(e))
            return RawDocument(content="", metadata={}, source_path=str(file_path))


# ---------------------------------------------------------------------------
# Markdown / Plain Text Loader
# ---------------------------------------------------------------------------
class TextLoader(DocumentLoader):
    """Plain text and Markdown file loader."""

    def can_load(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in {".txt", ".md", ".markdown", ".rst"}

    def load(self, file_path: Path) -> RawDocument:
        try:
            content = file_path.read_text(encoding="utf-8", errors="replace")

            metadata: dict = {
                "title": file_path.stem,
                "extension": file_path.suffix.lower(),
                "converter": "plaintext",
                "char_count": len(content),
            }

            return RawDocument(content=content, metadata=metadata, source_path=str(file_path))
        except Exception as e:
            logger.warning("text_load_failed", file=str(file_path), error=str(e))
            return RawDocument(content="", metadata={}, source_path=str(file_path))


# ---------------------------------------------------------------------------
# XML Loader
# ---------------------------------------------------------------------------
class XMLLoader(DocumentLoader):
    """XML file loader."""

    def can_load(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in {".xml", ".yaml", ".yml"}

    def load(self, file_path: Path) -> RawDocument:
        try:
            content = file_path.read_text(encoding="utf-8", errors="replace")

            metadata: dict = {
                "title": file_path.stem,
                "extension": file_path.suffix.lower(),
                "converter": "xml_text",
                "char_count": len(content),
            }

            return RawDocument(content=content, metadata=metadata, source_path=str(file_path))
        except Exception as e:
            logger.warning("xml_load_failed", file=str(file_path), error=str(e))
            return RawDocument(content="", metadata={}, source_path=str(file_path))
